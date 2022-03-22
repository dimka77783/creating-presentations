import time
import datetime
import openpyxl
from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
#  from auth_data import id_password
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

s = Service("E:\\РАБОЧИЕ ДОКУМЕНТЫ\\документы\\иницативы техсовет\\2021\\конструктор\\2TS bez EE\\chromedriver.exe")

path = "123.xlsx"  # имя файла
"""
n = input(str("Введите номер протокола: "))
date_prot = input('введите дату проведения протокола (гггг-мм-дд): ')
path_protokol = input('введите путь к протоколу: ')
date_now = datetime.date.today()

date_prot = date_prot.split('-')
day_count = datetime.date(int(date_prot[0]), int(date_prot[1]), int(date_prot[2]))
day_count = date_now - day_count
day_count = str(day_count)
day_count = (day_count.split()[0])
"""
wb_obj = openpyxl.load_workbook(path)  # Открываем файл
sheet_obj = wb_obj.active  # Выбираем активный лист таблицы(
m_row = sheet_obj.max_row

img_path = 'NLMK.png'  # загружаем фото
prs = Presentation()

url = "https://ideas.nlmk.com/"

options = webdriver.ChromeOptions()
options.add_argument("user-agent=Mozilla/5.0 (X11; Ubuntu; Linux x86_64; rv:84.0) Gecko/20100101 Firefox/84.0")
options.headless = True

driver = webdriver.Chrome(
    service=s,
    options=options
)

try:
    driver.get(url=url)
    time.sleep(10)
    """
    driver.find_element(By.CLASS_NAME, "auth-last-txt").click()
    time.sleep(5)

    
    email_input = driver.find_element(By.ID, "login")
    email_input.clear()
    email_input.send_keys("********")

    password_input = driver.find_element(By.ID, "password")
    password_input.clear()
    password_input.send_keys(id_password)
    time.sleep(3)
    password_input.send_keys(Keys.ENTER)
    time.sleep(6)
    """
    print('вход')
    driver.find_element(By.CLASS_NAME, "work").click()
    time.sleep(5)
    print('раб')
    driver.find_element(By.LINK_TEXT, "Все идеи").click()
    time.sleep(5)
    print('все')
    driver.find_element(By.LINK_TEXT, "ЦЖДТ").click()
    time.sleep(5)
    print('ЦЖДТ')
    """
    оперделяем количество страниц
    и заносим их в список 
    """
    
    my_ul2 = driver.find_elements(By.XPATH, "//ul[@class='b-page']")  # блок со страницам
    if len(my_ul2)==1:
        my_ul = driver.find_element(By.XPATH, "//ul[@class='b-page']")  # блок со страницам
        time.sleep(7)
        all_li = my_ul.find_elements(By.TAG_NAME, "li")  # оперделяем кол-во страниц
        time.sleep(7)
        link_list = []
        for li in all_li:
            y = li.text
            link_list.append(y)   
        rec = 1
        while rec == 1:
            for i in range(2, m_row + 1):
                cell_obj = sheet_obj.cell(row=i, column=1)  # ячейка с номером идей
                number_ideas = cell_obj.value
                cell_ob = sheet_obj.cell(row=i, column=10)  # В column= подставляем номер нужной колонки
                name_ideas = cell_ob.value
                cell_obj = sheet_obj.cell(row=i, column=10)  # В column= подставляем номер нужной колонки
                a = cell_obj.value
                cell_obj1 = sheet_obj.cell(row=i, column=11)  # В column= подставляем номер нужной колонки
                b = cell_obj1.value
                cell_obj2 = sheet_obj.cell(row=i, column=1)  # В column= подставляем номер нужной колонки
                z = cell_obj2.value
                cell_obj3 = sheet_obj.cell(row=i, column=3)  # В column= подставляем номер нужной колонки
                h = cell_obj3.value
                d = str(z) + ' ' + str(a)
                f = 'Руководитель проекта:' + str(h)
                cell_obj = sheet_obj.cell(row=i, column=12)  # В column= подставляем номер нужной колонки
                j = cell_obj.value
                cell_obj = sheet_obj.cell(row=i, column=13)  # В column= подставляем номер нужной колонки
                k = cell_obj.value
                cell_obj = sheet_obj.cell(row=i, column=6)  # В column= подставляем номер нужной колонки
                u = cell_obj.value
                
                for y in link_list:
                    y = y
                    ds = driver.find_elements(By.XPATH, " // *[ @ href = 'Offer.aspx?id=" + str(number_ideas) + "']")
                    time.sleep(7)
                    if len(ds) >= 1:
                        titleLayout = prs.slide_layouts[6]
                        slide = prs.slides.add_slide(titleLayout)

                        # рамещаем фото

                        left = Inches(0.5)
                        top = Inches(7)

                        pic = slide.shapes.add_picture(img_path, left, top)

                        # первая фигура №идей название и руководительпроекта
                        t1_left = Inches(0.5)  # расстояние от левого края
                        t1_top = Inches(0.5)  # расстояние от верха
                        t1_width = Inches(9)  # длинна
                        t1_height = Inches(2)  # ширина
                        txBox1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                        tf1 = txBox1.text_frame.paragraphs[0]
                        tf1.vertical_anchor = MSO_ANCHOR.TOP
                        tf1.word_wrap = True
                        tf1.margin_top = 0
                        tf1.horizontal_anchor = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                        run1 = tf1.add_run()
                        run1.text = d + '\n' + f
                        font = run1.font
                        font.name = 'Calibri'
                        font.size = Pt(24)  # размер шрифта
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.theme_color = MSO_THEME_COLOR_INDEX.DARK_1

                        # второй текст БЫЛО/ Текущее состояние

                        t2_left = Inches(0.5)  # расстояние от левого края
                        t2_top = Inches(1.5)  # расстояние от верха
                        t2_width = Inches(4)  # длинна
                        t2_height = Inches(1)  # ширина
                        txBox2 = slide.shapes.add_textbox(t2_left, t2_top, t2_width, t2_height)
                        tf2 = txBox2.text_frame.paragraphs[0]
                        tf2.horizontal_anchor = MSO_ANCHOR.TOP
                        tf2.word_wrap = True
                        tf2.margin_top = 0
                        tf2.horizontal_anchor = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        run2 = tf2.add_run()
                        run2.text = 'Было/Текущее состояние'
                        tf2.word_wrap = True
                        font = run2.font
                        font.name = 'Calibri'
                        font.size = Pt(24)  # размер шрифта
                        font.bold = True
                        font.italic = None
                        font.color.theme_color = MSO_THEME_COLOR_INDEX.DARK_1

                        # текс описание проблемы

                        textBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.5), Inches(1.0))
                        textFrame = textBox3.text_frame
                        textFrame.word_wrap = True
                        textParagraph = textFrame.add_paragraph()
                        textParagraph.text = 'Описание проблемы'

                        # текст самой проблемы

                        textBox4 = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(4.5), Inches(4.0))
                        textFrame = textBox4.text_frame
                        textFrame.word_wrap = True
                        textParagraph = textFrame.add_paragraph()
                        textParagraph.text = b

                        # возможные причины

                        textBox8 = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(4.5), Inches(4.0))
                        textFrame = textBox8.text_frame
                        textFrame.word_wrap = True
                        textParagraph = textFrame.add_paragraph()
                        textParagraph.text = 'Возможные причины'

                        # текс самой причины
                        """
                        textBox9 = slide.shapes.add_textbox(Inches(0.8), Inches(5.5),Inches(4.5), Inches(4.0))
                        textFrame = textBox9.text_frame
                        textFrame.word_wrap = True
                        textParagraph = textFrame.add_paragraph()
                        textParagraph.text = u

                        """
                        # текст Стало/ Будующее состояние

                        t5_left = Inches(5.2)  # расстояние от левого края
                        t5_top = Inches(1.5)  # расстояние от верха
                        t5_width = Inches(4)  # длинна
                        t5_height = Inches(1)  # ширина
                        txBox5 = slide.shapes.add_textbox(t5_left, t5_top, t5_width, t5_height)
                        tf5 = txBox5.text_frame.paragraphs[0]
                        tf5.horizontal_anchor = MSO_ANCHOR.TOP
                        tf5.word_wrap = True
                        tf5.margin_top = 0
                        tf5.horizontal_anchor = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        run5 = tf5.add_run()
                        run5.text = 'Стало / Будущее состояние'
                        tf5.word_wrap = True
                        font = run5.font
                        font.name = 'Calibri'
                        font.size = Pt(24)  # размер шрифта
                        font.bold = True
                        font.italic = None  # cause value to be inherited from theme
                        font.color.theme_color = MSO_THEME_COLOR_INDEX.DARK_1
                        # текс описание проблемы

                        textBox6 = slide.shapes.add_textbox(Inches(5.3), Inches(1.6), Inches(4), Inches(1.0))
                        textFrame = textBox6.text_frame
                        textFrame.word_wrap = True
                        textParagraph = textFrame.add_paragraph()
                        textParagraph.text = 'Предлагаемые изменения'

                        # текст самого решения

                        textBox7 = slide.shapes.add_textbox(Inches(5.6), Inches(2), Inches(4), Inches(1.0))
                        textFrame = textBox7.text_frame
                        textFrame.word_wrap = True
                        textParagraph = textFrame.add_paragraph()
                        textParagraph.text = str(j) + '\n' + '\n''Ожидаемый результат' + '\n' + '  ' + str(k)
                        driver.find_element(By.XPATH, " // *[ @ href = 'Offer.aspx?id=" + str(number_ideas) + "']").click()
                        time.sleep(8)
                        elems = driver.find_elements(By.XPATH,"//a[contains(@id,'HyperLinkName')]")
                        for elem in elems:
                            past = elem.get_attribute("href")
                            #  print(elem.get_attribute("href"))
                            di = past.split('=')
                            number_file = di[2]
                            #  print(number_file)
                            driver.find_element(By.XPATH, " // *[ @ href = 'DownloadFile.aspx?obj=Offer&id=" + str(
                                number_file) + "']").click()
                            time.sleep(2)
                            get = elem.get_attribute("text")
                            #  print(type(get))
                            print(get)
                            get = get.replace('~','_')
                            time.sleep(2)

                            img_path2 = "C:\\Users\\odinokov_da\\Downloads\\"+get+""
                            left = Inches(0.5)
                            top = Inches(5)

                            pic2 = slide.shapes.add_picture(img_path2, left, top, Inches(3))

                        driver.find_element(By.LINK_TEXT, "Все идеи").click()
                        time.sleep(5)
                        prs.save('123.pptx')
                        break

                    else:
                        time.sleep(5)
                        rec = 0
                        w = int(y) + 1
                        tt = int(y)+1

                        print(f"переходим на страницу: {w}")
                        driver.find_element(By.LINK_TEXT,  "" + str(tt) + "").click()
                        time.sleep(7)
        prs.save('123.pptx')
    else:
        for i in range(2, m_row + 1):
            cell_obj = sheet_obj.cell(row=i, column=1)  # ячейка с номером идей
            number_ideas = cell_obj.value
            cell_ob = sheet_obj.cell(row=i, column=10)  # В column= подставляем номер нужной колонки
            name_ideas = cell_ob.value
            cell_obj = sheet_obj.cell(row=i, column=10)  # В column= подставляем номер нужной колонки
            a = cell_obj.value
            cell_obj1 = sheet_obj.cell(row=i, column=11)  # В column= подставляем номер нужной колонки
            b = cell_obj1.value
            cell_obj2 = sheet_obj.cell(row=i, column=1)  # В column= подставляем номер нужной колонки
            z = cell_obj2.value
            cell_obj3 = sheet_obj.cell(row=i, column=3)  # В column= подставляем номер нужной колонки
            h = cell_obj3.value
            d = str(z) + ' ' + str(a)
            f = 'Руководитель проекта:' + str(h)
            cell_obj = sheet_obj.cell(row=i, column=12)  # В column= подставляем номер нужной колонки
            j = cell_obj.value
            cell_obj = sheet_obj.cell(row=i, column=13)  # В column= подставляем номер нужной колонки
            k = cell_obj.value
            cell_obj = sheet_obj.cell(row=i, column=6)  # В column= подставляем номер нужной колонки
            u = cell_obj.value
                
            for y in link_list:
                y = y
                ds = driver.find_elements(By.XPATH, " // *[ @ href = 'Offer.aspx?id=" + str(number_ideas) + "']")
                time.sleep(7)
                if len(ds) >= 1:
                    titleLayout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(titleLayout)

                    # рамещаем фото

                    left = Inches(0.5)
                    top = Inches(7)

                    pic = slide.shapes.add_picture(img_path, left, top)

                    # первая фигура №идей название и руководительпроекта
                    t1_left = Inches(0.5)  # расстояние от левого края
                    t1_top = Inches(0.5)  # расстояние от верха
                    t1_width = Inches(9)  # длинна
                    t1_height = Inches(2)  # ширина
                    txBox1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                    tf1 = txBox1.text_frame.paragraphs[0]
                    tf1.vertical_anchor = MSO_ANCHOR.TOP
                    tf1.word_wrap = True
                    tf1.margin_top = 0
                    tf1.horizontal_anchor = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    run1 = tf1.add_run()
                    run1.text = d + '\n' + f
                    font = run1.font
                    font.name = 'Calibri'
                    font.size = Pt(24)  # размер шрифта
                    font.bold = True
                    font.italic = None  # cause value to be inherited from theme
                    font.color.theme_color = MSO_THEME_COLOR_INDEX.DARK_1

                    # второй текст БЫЛО/ Текущее состояние

                    t2_left = Inches(0.5)  # расстояние от левого края
                    t2_top = Inches(1.5)  # расстояние от верха
                    t2_width = Inches(4)  # длинна
                    t2_height = Inches(1)  # ширина
                    txBox2 = slide.shapes.add_textbox(t2_left, t2_top, t2_width, t2_height)
                    tf2 = txBox2.text_frame.paragraphs[0]
                    tf2.horizontal_anchor = MSO_ANCHOR.TOP
                    tf2.word_wrap = True
                    tf2.margin_top = 0
                    tf2.horizontal_anchor = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    run2 = tf2.add_run()
                    run2.text = 'Было/Текущее состояние'
                    tf2.word_wrap = True
                    font = run2.font
                    font.name = 'Calibri'
                    font.size = Pt(24)  # размер шрифта
                    font.bold = True
                    font.italic = None
                    font.color.theme_color = MSO_THEME_COLOR_INDEX.DARK_1

                    # текс описание проблемы

                    textBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.5), Inches(1.0))
                    textFrame = textBox3.text_frame
                    textFrame.word_wrap = True
                    textParagraph = textFrame.add_paragraph()
                    textParagraph.text = 'Описание проблемы'

                    # текст самой проблемы

                    textBox4 = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(4.5), Inches(4.0))
                    textFrame = textBox4.text_frame
                    textFrame.word_wrap = True
                    textParagraph = textFrame.add_paragraph()
                    textParagraph.text = b

                    # возможные причины

                    textBox8 = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(4.5), Inches(4.0))
                    textFrame = textBox8.text_frame
                    textFrame.word_wrap = True
                    textParagraph = textFrame.add_paragraph()
                    textParagraph.text = 'Возможные причины'

                    # текс самой причины
                    """
                    textBox9 = slide.shapes.add_textbox(Inches(0.8), Inches(5.5),Inches(4.5), Inches(4.0))
                    textFrame = textBox9.text_frame
                    textFrame.word_wrap = True
                    textParagraph = textFrame.add_paragraph()
                    textParagraph.text = u

                    """
                    # текст Стало/ Будующее состояние

                    t5_left = Inches(5.2)  # расстояние от левого края
                    t5_top = Inches(1.5)  # расстояние от верха
                    t5_width = Inches(4)  # длинна
                    t5_height = Inches(1)  # ширина
                    txBox5 = slide.shapes.add_textbox(t5_left, t5_top, t5_width, t5_height)
                    tf5 = txBox5.text_frame.paragraphs[0]
                    tf5.horizontal_anchor = MSO_ANCHOR.TOP
                    tf5.word_wrap = True
                    tf5.margin_top = 0
                    tf5.horizontal_anchor = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    run5 = tf5.add_run()
                    run5.text = 'Стало / Будущее состояние'
                    tf5.word_wrap = True
                    font = run5.font
                    font.name = 'Calibri'
                    font.size = Pt(24)  # размер шрифта
                    font.bold = True
                    font.italic = None  # cause value to be inherited from theme
                    font.color.theme_color = MSO_THEME_COLOR_INDEX.DARK_1
                    # текс описание проблемы

                    textBox6 = slide.shapes.add_textbox(Inches(5.3), Inches(1.6), Inches(4), Inches(1.0))
                    textFrame = textBox6.text_frame
                    textFrame.word_wrap = True
                    textParagraph = textFrame.add_paragraph()
                    textParagraph.text = 'Предлагаемые изменения'

                    # текст самого решения

                    textBox7 = slide.shapes.add_textbox(Inches(5.6), Inches(2), Inches(4), Inches(1.0))
                    textFrame = textBox7.text_frame
                    textFrame.word_wrap = True
                    textParagraph = textFrame.add_paragraph()
                    textParagraph.text = str(j) + '\n' + '\n''Ожидаемый результат' + '\n' + '  ' + str(k)
                    driver.find_element(By.XPATH, " // *[ @ href = 'Offer.aspx?id=" + str(number_ideas) + "']").click()
                    time.sleep(8)
                    elems = driver.find_elements(By.XPATH,"//a[contains(@id,'HyperLinkName')]")
                    for elem in elems:
                        past = elem.get_attribute("href")
                        #  print(elem.get_attribute("href"))
                        di = past.split('=')
                        number_file = di[2]
                        #  print(number_file)
                        driver.find_element(By.XPATH, " // *[ @ href = 'DownloadFile.aspx?obj=Offer&id=" + str(
                            number_file) + "']").click()
                        get = elem.get_attribute("text")
                        # print(type(get))
                        print(get)

                        img_path2 = "C:\\Users\\odinokov_da\\Downloads\\"+get+""
                        left = Inches(0.5)
                        top = Inches(5)

                        pic2 = slide.shapes.add_picture(img_path2, left, top, Inches(3))

                    driver.find_element(By.LINK_TEXT, "Все идеи").click()
                    time.sleep(5)                       
    prs.save('123.pptx')
         
except Exception as ex:
    print(ex)
finally:
    driver.close()
    driver.quit()
